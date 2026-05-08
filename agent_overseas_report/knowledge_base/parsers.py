"""Local document text extraction utilities for the knowledge base."""

from __future__ import annotations

from dataclasses import dataclass, field
from pathlib import Path
from typing import Any

SUPPORTED_EXTENSIONS = {
    ".pdf": "pdf",
    ".docx": "word",
    ".doc": "word",
    ".xlsx": "excel",
    ".xlsm": "excel",
    ".xls": "excel",
    ".pptx": "ppt",
    ".ppt": "ppt",
    ".txt": "txt",
    ".md": "markdown",
    ".markdown": "markdown",
}


@dataclass(slots=True)
class ParsedTextBlock:
    """One parsed text block with document-position metadata."""

    text: str
    page_number: int | None = None
    sheet_name: str | None = None
    slide_number: int | None = None
    metadata: dict[str, Any] = field(default_factory=dict)


def identify_file_type(file_name: str, content_type: str | None = None) -> str:
    """Identify a knowledge file type from extension and upload content type."""

    suffix = Path(file_name).suffix.lower()
    if suffix in SUPPORTED_EXTENSIONS:
        return SUPPORTED_EXTENSIONS[suffix]
    if content_type:
        normalized = content_type.lower()
        if "pdf" in normalized:
            return "pdf"
        if "word" in normalized or "officedocument.wordprocessingml" in normalized:
            return "word"
        if "excel" in normalized or "spreadsheetml" in normalized:
            return "excel"
        if "powerpoint" in normalized or "presentationml" in normalized:
            return "ppt"
        if "markdown" in normalized:
            return "markdown"
        if normalized.startswith("text/"):
            return "txt"
    return "unknown"


def parse_document(file_path: Path, file_type: str) -> list[ParsedTextBlock]:
    """Extract text blocks from a supported local document."""

    if file_type == "pdf":
        return _parse_pdf(file_path)
    if file_type == "word":
        return _parse_word(file_path)
    if file_type == "excel":
        return _parse_excel(file_path)
    if file_type == "ppt":
        return _parse_ppt(file_path)
    if file_type in {"txt", "markdown"}:
        return _parse_text(file_path, file_type)
    raise ValueError(f"Unsupported file type: {file_type}")


def _parse_pdf(file_path: Path) -> list[ParsedTextBlock]:
    from pypdf import PdfReader

    reader = PdfReader(str(file_path))
    blocks: list[ParsedTextBlock] = []
    for index, page in enumerate(reader.pages, start=1):
        text = (page.extract_text() or "").strip()
        if text:
            blocks.append(ParsedTextBlock(text=text, page_number=index, metadata={"parser": "pypdf"}))
    return blocks


def _parse_word(file_path: Path) -> list[ParsedTextBlock]:
    from docx import Document

    document = Document(str(file_path))
    paragraphs = [paragraph.text.strip() for paragraph in document.paragraphs if paragraph.text.strip()]
    table_lines: list[str] = []
    for table_index, table in enumerate(document.tables, start=1):
        for row_index, row in enumerate(table.rows, start=1):
            values = [cell.text.strip() for cell in row.cells if cell.text.strip()]
            if values:
                table_lines.append(f"表格{table_index} 行{row_index}: " + " | ".join(values))
    text = "\n".join(paragraphs + table_lines).strip()
    return [ParsedTextBlock(text=text, metadata={"parser": "python-docx"})] if text else []


def _parse_excel(file_path: Path) -> list[ParsedTextBlock]:
    from openpyxl import load_workbook

    workbook = load_workbook(str(file_path), data_only=True, read_only=True)
    blocks: list[ParsedTextBlock] = []
    for worksheet in workbook.worksheets:
        rows: list[str] = []
        for row in worksheet.iter_rows(values_only=True):
            values = [str(value).strip() for value in row if value is not None and str(value).strip()]
            if values:
                rows.append(" | ".join(values))
        text = "\n".join(rows).strip()
        if text:
            blocks.append(ParsedTextBlock(text=text, sheet_name=worksheet.title, metadata={"parser": "openpyxl"}))
    workbook.close()
    return blocks


def _parse_ppt(file_path: Path) -> list[ParsedTextBlock]:
    from pptx import Presentation

    presentation = Presentation(str(file_path))
    blocks: list[ParsedTextBlock] = []
    for slide_number, slide in enumerate(presentation.slides, start=1):
        lines: list[str] = []
        for shape in slide.shapes:
            text = getattr(shape, "text", "")
            if text and text.strip():
                lines.append(text.strip())
        text = "\n".join(lines).strip()
        if text:
            blocks.append(ParsedTextBlock(text=text, slide_number=slide_number, metadata={"parser": "python-pptx"}))
    return blocks


def _parse_text(file_path: Path, file_type: str) -> list[ParsedTextBlock]:
    raw = file_path.read_bytes()
    try:
        text = raw.decode("utf-8")
        encoding = "utf-8"
    except UnicodeDecodeError:
        text = raw.decode("gb18030", errors="replace")
        encoding = "gb18030"
    text = text.strip()
    return [ParsedTextBlock(text=text, metadata={"parser": file_type, "encoding": encoding})] if text else []


def split_text_blocks(blocks: list[ParsedTextBlock], *, max_chars: int = 1200, overlap: int = 120) -> list[ParsedTextBlock]:
    """Split parsed blocks into chunk-sized blocks while preserving location metadata."""

    chunks: list[ParsedTextBlock] = []
    for block in blocks:
        text = "\n".join(line.strip() for line in block.text.splitlines() if line.strip())
        if not text:
            continue
        if len(text) <= max_chars:
            chunks.append(block)
            continue
        start = 0
        while start < len(text):
            end = min(start + max_chars, len(text))
            chunk_text = text[start:end].strip()
            if chunk_text:
                chunks.append(
                    ParsedTextBlock(
                        text=chunk_text,
                        page_number=block.page_number,
                        sheet_name=block.sheet_name,
                        slide_number=block.slide_number,
                        metadata={**block.metadata, "source_offset": start},
                    )
                )
            if end == len(text):
                break
            start = max(end - overlap, start + 1)
    return chunks


def estimate_token_count(text: str) -> int:
    """Estimate tokens for storage until a tokenizer is introduced."""

    ascii_words = len([part for part in text.split() if part])
    cjk_chars = sum(1 for char in text if "\u4e00" <= char <= "\u9fff")
    return max(1, ascii_words + cjk_chars)
